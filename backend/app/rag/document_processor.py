"""
Document Processor avec extraction hybride (Texte + Images OCR).

CORRECTION : Estimation correcte du nombre de pages pour DOCX/PPTX/XLSX.
NORME : Tous les enums en MAJUSCULES.
"""
import logging
import io
import math
from pathlib import Path
from typing import Dict, List, Optional, Union, Any

# PDF
try:
    from pypdf import PdfReader
except ImportError:
    from PyPDF2 import PdfReader

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

# DOCX
try:
    from docx import Document as DocxDocument
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PPTX
try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# XLSX
try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

# Images
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

# OCR Client
from app.rag.ocr_processor import MistralOCRClient
from app.core.config import settings

logger = logging.getLogger(__name__)


# =============================================================================
# CONSTANTES
# =============================================================================

class ExtractionMethod:
    """Méthodes d'extraction - MAJUSCULES."""
    TEXT = "TEXT"
    OCR = "OCR"
    HYBRID = "HYBRID"
    FALLBACK = "FALLBACK"


# Estimation du nombre de caractères par page (moyenne)
# Source : Une page A4 standard contient environ 2000-3000 caractères
CHARS_PER_PAGE_ESTIMATE = 2500

# Estimation pour les paragraphes (environ 15-20 paragraphes par page)
PARAGRAPHS_PER_PAGE_ESTIMATE = 15


class DocumentProcessor:
    """
    Processeur de documents avec extraction hybride.
    
    CORRECTION : Estimation correcte du nombre de pages.
    NORME : extraction_method toujours en MAJUSCULES.
    """
    
    def __init__(self, ocr_client: Optional[MistralOCRClient] = None):
        """
        Initialiser le processeur.
        
        Args:
            ocr_client: Client OCR Mistral (créé automatiquement si non fourni)
        """
        if ocr_client:
            self.ocr_client = ocr_client
        else:
            self.ocr_client = MistralOCRClient(
                api_key=settings.MISTRAL_API_KEY,
                model="mistral-ocr-latest"
            )
        
        logger.info("DocumentProcessor initialized with OCR support")
    
    # =========================================================================
    # MÉTHODE PRINCIPALE
    # =========================================================================
    
    def process_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Traiter un document et extraire tout le contenu (texte + images).
        
        Args:
            file_path: Chemin vers le fichier
            
        Returns:
            Dict avec:
                - text: Texte combiné (texte + OCR images)
                - pages: Liste des pages avec détails
                - page_count: Nombre de pages (estimé si nécessaire)
                - total_images: Nombre d'images traitées
                - extraction_method: "TEXT", "OCR", ou "HYBRID" (MAJUSCULES)
                - file_type: Type de fichier
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        # Router vers la bonne méthode d'extraction
        extractors = {
            ".pdf": self.extract_pdf_hybrid,
            ".docx": self.extract_docx_hybrid,
            ".doc": self.extract_docx_hybrid,
            ".pptx": self.extract_pptx_hybrid,
            ".ppt": self.extract_pptx_hybrid,
            ".xlsx": self.extract_xlsx,
            ".xls": self.extract_xlsx,
            ".txt": self.extract_text,
            ".md": self.extract_text,
            ".rtf": self.extract_rtf,
            ".png": self.extract_image,
            ".jpg": self.extract_image,
            ".jpeg": self.extract_image,
            ".webp": self.extract_image,
            ".gif": self.extract_image,
        }
        
        extractor = extractors.get(extension)
        
        if not extractor:
            raise ValueError(f"Format non supporté: {extension}")
        
        result = extractor(file_path)
        result["file_type"] = extension.lstrip(".")
        
        # S'assurer que page_count est présent
        if "page_count" not in result:
            result["page_count"] = len(result.get("pages", [1]))
        
        logger.info(
            f"Processed {file_path.name}: {len(result['text'])} chars, "
            f"{result['page_count']} pages, {result['total_images']} images, "
            f"method={result['extraction_method']}"
        )
        
        return result
    
    # =========================================================================
    # UTILITAIRE : ESTIMATION DU NOMBRE DE PAGES
    # =========================================================================
    
    @staticmethod
    def estimate_page_count(
        text_length: int = 0,
        paragraph_count: int = 0,
        method: str = "chars"
    ) -> int:
        """
        Estimer le nombre de pages d'un document.
        
        Args:
            text_length: Longueur du texte en caractères
            paragraph_count: Nombre de paragraphes
            method: Méthode d'estimation ("chars", "paragraphs", "combined")
            
        Returns:
            Nombre de pages estimé (minimum 1)
        """
        if method == "chars":
            pages = math.ceil(text_length / CHARS_PER_PAGE_ESTIMATE)
        elif method == "paragraphs":
            pages = math.ceil(paragraph_count / PARAGRAPHS_PER_PAGE_ESTIMATE)
        elif method == "combined":
            # Moyenne des deux méthodes
            pages_by_chars = math.ceil(text_length / CHARS_PER_PAGE_ESTIMATE)
            pages_by_paras = math.ceil(paragraph_count / PARAGRAPHS_PER_PAGE_ESTIMATE)
            pages = math.ceil((pages_by_chars + pages_by_paras) / 2)
        else:
            pages = 1
        
        return max(1, pages)
    
    # =========================================================================
    # PDF - EXTRACTION HYBRIDE
    # =========================================================================
    
    def extract_pdf_hybrid(self, pdf_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extraction hybride PDF : texte natif + OCR pour images/pages scannées.
        
        Returns:
            Dict avec page_count exact (PDF a des pages natives)
        """
        pdf_path = Path(pdf_path)
        
        result = {
            "text": "",
            "pages": [],
            "page_count": 0,
            "total_images": 0,
            "extraction_method": ExtractionMethod.TEXT
        }
        
        try:
            reader = PdfReader(str(pdf_path))
            total_pages = len(reader.pages)
            result["page_count"] = total_pages  # PDF a des pages exactes
            
            # Première passe : extraire texte natif
            native_texts = []
            for page in reader.pages:
                text = page.extract_text() or ""
                native_texts.append(text.strip())
            
            # Calculer si PDF est principalement scanné
            total_native_chars = sum(len(t) for t in native_texts)
            avg_chars_per_page = total_native_chars / max(total_pages, 1)
            
            # Seuil : si moins de 50 chars/page en moyenne → PDF scanné
            is_scanned_pdf = avg_chars_per_page < 50
            
            if is_scanned_pdf:
                # PDF SCANNÉ : OCR sur tout le document
                logger.info(f"PDF scanné détecté ({avg_chars_per_page:.0f} chars/page), OCR complet")
                result = self._ocr_entire_pdf(pdf_path)
                result["extraction_method"] = ExtractionMethod.OCR
                result["page_count"] = total_pages  # Garder le vrai nombre de pages
            
            else:
                # PDF STANDARD : texte natif + OCR images intégrées
                logger.info(f"PDF standard ({avg_chars_per_page:.0f} chars/page), extraction hybride")
                result = self._extract_pdf_with_embedded_images(pdf_path, reader, native_texts)
                result["page_count"] = total_pages  # Garder le vrai nombre de pages
                
                # Déterminer méthode
                if result["total_images"] > 0:
                    result["extraction_method"] = ExtractionMethod.HYBRID
                else:
                    result["extraction_method"] = ExtractionMethod.TEXT
        
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            # Fallback : OCR complet
            logger.info("Fallback to full OCR")
            result = self._ocr_entire_pdf(pdf_path)
            result["extraction_method"] = ExtractionMethod.OCR
        
        return result
    
    def _ocr_entire_pdf(self, pdf_path: Path) -> Dict[str, Any]:
        """OCR complet d'un PDF (pour PDF scannés)."""
        result = {
            "text": "",
            "pages": [],
            "page_count": 1,
            "total_images": 0,
            "extraction_method": ExtractionMethod.OCR
        }
        
        try:
            ocr_text = self.ocr_client.extract_text_from_pdf(str(pdf_path))
            
            result["text"] = ocr_text
            result["pages"] = [{"page_num": 1, "text": ocr_text, "method": "OCR"}]
            
            # Compter les pages dans le résultat OCR
            page_count = ocr_text.count("=== Page ")
            if page_count == 0:
                # Estimer à partir du texte
                page_count = self.estimate_page_count(len(ocr_text))
            
            result["page_count"] = page_count
            result["total_images"] = page_count  # Chaque page scannée = 1 image
        
        except Exception as e:
            logger.error(f"Full PDF OCR failed: {e}")
            result["text"] = f"[OCR Error: {str(e)}]"
        
        return result
    
    def _extract_pdf_with_embedded_images(
        self, 
        pdf_path: Path, 
        reader: PdfReader,
        native_texts: List[str]
    ) -> Dict[str, Any]:
        """Extraire PDF standard avec images intégrées."""
        result = {
            "text": "",
            "pages": [],
            "page_count": len(reader.pages),
            "total_images": 0,
            "extraction_method": ExtractionMethod.HYBRID
        }
        
        for page_num, (page, native_text) in enumerate(zip(reader.pages, native_texts), 1):
            page_data = {
                "page_num": page_num,
                "native_text": native_text,
                "image_texts": [],
                "combined_text": native_text
            }
            
            # Extraire images de la page
            try:
                if "/XObject" in page.get("/Resources", {}):
                    xobject = page["/Resources"]["/XObject"].get_object()
                    
                    for obj_name in xobject:
                        obj = xobject[obj_name]
                        
                        if obj.get("/Subtype") == "/Image":
                            try:
                                image_data = self._extract_image_from_pdf_object(obj)
                                if image_data:
                                    ocr_text = self.ocr_client.extract_text_from_image(image_data)
                                    
                                    if ocr_text and not ocr_text.startswith("["):
                                        page_data["image_texts"].append(ocr_text)
                                        result["total_images"] += 1
                            except Exception as img_err:
                                logger.warning(f"Failed to OCR image {obj_name}: {img_err}")
            
            except Exception as e:
                logger.warning(f"Failed to extract images from page {page_num}: {e}")
            
            # Combiner texte natif + OCR images
            combined = native_text
            for idx, img_text in enumerate(page_data["image_texts"], 1):
                combined += f"\n\n[Image {idx}]\n{img_text}"
            
            page_data["combined_text"] = combined
            result["pages"].append(page_data)
            result["text"] += f"\n\n--- Page {page_num} ---\n{combined}"
        
        return result
    
    def _extract_image_from_pdf_object(self, pdf_image_obj) -> Optional[Image.Image]:
        """Extraire une image PIL depuis un objet PDF XObject."""
        if not PIL_AVAILABLE:
            return None
        
        try:
            data = pdf_image_obj.get_data()
            image = Image.open(io.BytesIO(data))
            return image
        except Exception as e:
            logger.debug(f"Could not extract image from PDF object: {e}")
            return None
    
    # =========================================================================
    # DOCX - EXTRACTION HYBRIDE AVEC ESTIMATION PAGES
    # =========================================================================
    
    def extract_docx_hybrid(self, docx_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extraction hybride DOCX : texte + OCR images.
        
        CORRECTION : Estimation du nombre de pages basée sur le contenu.
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx requis: pip install python-docx")
        
        docx_path = Path(docx_path)
        
        result = {
            "text": "",
            "pages": [],
            "page_count": 1,
            "total_images": 0,
            "extraction_method": ExtractionMethod.TEXT
        }
        
        try:
            doc = DocxDocument(str(docx_path))
            
            # 1. Extraire texte des paragraphes
            paragraphs_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs_text.append(para.text)
            
            # 2. Extraire texte des tableaux
            tables_text = []
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_content.append(row_text)
                if table_content:
                    tables_text.append("\n".join(table_content))
            
            # 3. Extraire et OCR les images
            image_texts = []
            
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_data = rel.target_part.blob
                        image = Image.open(io.BytesIO(image_data))
                        
                        ocr_text = self.ocr_client.extract_text_from_image(image)
                        
                        if ocr_text and not ocr_text.startswith("["):
                            image_texts.append(ocr_text)
                            result["total_images"] += 1
                    
                    except Exception as img_err:
                        logger.warning(f"Failed to OCR DOCX image: {img_err}")
            
            # 4. Combiner tout
            all_text = "\n\n".join(paragraphs_text)
            
            if tables_text:
                all_text += "\n\n--- Tableaux ---\n" + "\n\n".join(tables_text)
            
            for idx, img_text in enumerate(image_texts, 1):
                all_text += f"\n\n[Image {idx}]\n{img_text}"
            
            result["text"] = all_text
            
            # 5. CORRECTION : Estimer le nombre de pages
            paragraph_count = len(paragraphs_text)
            text_length = len(all_text)
            
            # Utiliser la méthode combinée pour une meilleure estimation
            estimated_pages = self.estimate_page_count(
                text_length=text_length,
                paragraph_count=paragraph_count,
                method="combined"
            )
            
            result["page_count"] = estimated_pages
            
            result["pages"] = [{
                "page_num": i + 1,
                "estimated": True
            } for i in range(estimated_pages)]
            
            # Ajouter des infos détaillées dans la première "page"
            result["pages"][0].update({
                "paragraphs": paragraph_count,
                "tables": len(tables_text),
                "images": len(image_texts),
                "text_length": text_length
            })
            
            # Déterminer méthode
            if result["total_images"] > 0:
                result["extraction_method"] = ExtractionMethod.HYBRID
            else:
                result["extraction_method"] = ExtractionMethod.TEXT
            
            logger.info(
                f"DOCX extracted: {text_length} chars, {paragraph_count} paragraphs, "
                f"estimated {estimated_pages} pages, {result['total_images']} images"
            )
        
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            result["text"] = f"[Extraction Error: {str(e)}]"
        
        return result
    
    # =========================================================================
    # PPTX - EXTRACTION HYBRIDE (slides = pages)
    # =========================================================================
    
    def extract_pptx_hybrid(self, pptx_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extraction hybride PPTX : texte slides + OCR images.
        
        PPTX : 1 slide = 1 page (pas d'estimation nécessaire)
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx requis: pip install python-pptx")
        
        pptx_path = Path(pptx_path)
        
        result = {
            "text": "",
            "pages": [],
            "page_count": 0,
            "total_images": 0,
            "extraction_method": ExtractionMethod.TEXT
        }
        
        try:
            prs = Presentation(str(pptx_path))
            slide_count = len(prs.slides)
            result["page_count"] = slide_count  # 1 slide = 1 page
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "page_num": slide_num,
                    "slide_num": slide_num,
                    "text_content": [],
                    "image_texts": []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["text_content"].append(shape.text)
                    
                    if shape.has_table:
                        table_text = []
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                table_text.append(row_text)
                        if table_text:
                            slide_data["text_content"].append("[Tableau]\n" + "\n".join(table_text))
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_blob = shape.image.blob
                            image = Image.open(io.BytesIO(image_blob))
                            
                            ocr_text = self.ocr_client.extract_text_from_image(image)
                            
                            if ocr_text and not ocr_text.startswith("["):
                                slide_data["image_texts"].append(ocr_text)
                                result["total_images"] += 1
                        
                        except Exception as img_err:
                            logger.warning(f"Failed to OCR PPTX image: {img_err}")
                
                slide_text = "\n".join(slide_data["text_content"])
                
                for idx, img_text in enumerate(slide_data["image_texts"], 1):
                    slide_text += f"\n\n[Image {idx}]\n{img_text}"
                
                slide_data["combined_text"] = slide_text
                result["pages"].append(slide_data)
                result["text"] += f"\n\n--- Slide {slide_num} ---\n{slide_text}"
            
            # Déterminer méthode
            if result["total_images"] > 0:
                result["extraction_method"] = ExtractionMethod.HYBRID
            else:
                result["extraction_method"] = ExtractionMethod.TEXT
        
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            result["text"] = f"[Extraction Error: {str(e)}]"
        
        return result
    
    # =========================================================================
    # XLSX - EXTRACTION (feuilles = pages)
    # =========================================================================
    
    def extract_xlsx(self, xlsx_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extraction XLSX : texte des cellules.
        
        XLSX : 1 feuille = 1 page
        """
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl requis: pip install openpyxl")
        
        xlsx_path = Path(xlsx_path)
        
        result = {
            "text": "",
            "pages": [],
            "page_count": 0,
            "total_images": 0,
            "extraction_method": ExtractionMethod.TEXT
        }
        
        try:
            wb = load_workbook(str(xlsx_path), data_only=True)
            sheet_count = len(wb.sheetnames)
            result["page_count"] = sheet_count  # 1 feuille = 1 page
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                sheet_data = {
                    "page_num": wb.sheetnames.index(sheet_name) + 1,
                    "sheet_name": sheet_name,
                    "rows": []
                }
                
                sheet_text = f"[Feuille: {sheet_name}]\n"
                
                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    
                    if row_values:
                        row_text = " | ".join(row_values)
                        sheet_data["rows"].append(row_text)
                        sheet_text += row_text + "\n"
                
                result["pages"].append(sheet_data)
                result["text"] += sheet_text + "\n"
        
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            result["text"] = f"[Extraction Error: {str(e)}]"
        
        return result
    
    # =========================================================================
    # IMAGES PURES - OCR DIRECT
    # =========================================================================
    
    def extract_image(self, image_path: Union[str, Path]) -> Dict[str, Any]:
        """Extraction image pure : OCR direct."""
        image_path = Path(image_path)
        
        result = {
            "text": "",
            "pages": [],
            "page_count": 1,  # 1 image = 1 page
            "total_images": 1,
            "extraction_method": ExtractionMethod.OCR
        }
        
        try:
            ocr_text = self.ocr_client.extract_text_from_image(str(image_path))
            
            result["text"] = ocr_text
            result["pages"] = [{"page_num": 1, "text": ocr_text, "method": "OCR"}]
        
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            result["text"] = f"[OCR Error: {str(e)}]"
        
        return result
    
    # =========================================================================
    # FICHIERS TEXTE
    # =========================================================================
    
    def extract_text(self, text_path: Union[str, Path]) -> Dict[str, Any]:
        """Extraction fichier texte (TXT, MD) avec estimation pages."""
        text_path = Path(text_path)
        
        result = {
            "text": "",
            "pages": [],
            "page_count": 1,
            "total_images": 0,
            "extraction_method": ExtractionMethod.TEXT
        }
        
        try:
            encodings = ["utf-8", "latin-1", "cp1252"]
            
            for encoding in encodings:
                try:
                    with open(text_path, "r", encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                content = "[Encoding Error: Unable to decode file]"
            
            result["text"] = content
            
            # Estimer le nombre de pages
            result["page_count"] = self.estimate_page_count(len(content))
            result["pages"] = [{"page_num": 1, "text": content}]
        
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            result["text"] = f"[Extraction Error: {str(e)}]"
        
        return result
    
    # =========================================================================
    # RTF
    # =========================================================================
    
    def extract_rtf(self, rtf_path: Union[str, Path]) -> Dict[str, Any]:
        """Extraction fichier RTF avec estimation pages."""
        rtf_path = Path(rtf_path)
        
        result = {
            "text": "",
            "pages": [],
            "page_count": 1,
            "total_images": 0,
            "extraction_method": ExtractionMethod.TEXT
        }
        
        try:
            try:
                from striprtf.striprtf import rtf_to_text
                
                with open(rtf_path, "r", encoding="utf-8", errors="ignore") as f:
                    rtf_content = f.read()
                
                content = rtf_to_text(rtf_content)
            
            except ImportError:
                with open(rtf_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                
                import re
                content = re.sub(r"\\[a-z]+\d*\s?", "", content)
                content = re.sub(r"[{}]", "", content)
            
            result["text"] = content
            
            # Estimer le nombre de pages
            result["page_count"] = self.estimate_page_count(len(content))
            result["pages"] = [{"page_num": 1, "text": content}]
        
        except Exception as e:
            logger.error(f"RTF extraction failed: {e}")
            result["text"] = f"[Extraction Error: {str(e)}]"
        
        return result


# =============================================================================
# FONCTION HELPER
# =============================================================================

def process_document_with_ocr(
    file_path: str,
    ocr_client: Optional[MistralOCRClient] = None
) -> Dict[str, Any]:
    """
    Fonction helper pour traiter un document avec OCR.
    
    Returns:
        Dict avec extraction_method en MAJUSCULES et page_count estimé
    """
    processor = DocumentProcessor(ocr_client=ocr_client)
    return processor.process_document(file_path)
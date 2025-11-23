"""
Processing worker tasks for document text extraction with Hybrid OCR support.

CORRECTIONS :
- Nettoyage des caractères NULL (\u0000) avant stockage PostgreSQL
- Estimation correcte du nombre de pages pour DOCX
- Enums en MAJUSCULES
"""
import logging
import re
from typing import Optional, Dict, Any
from datetime import datetime
from pathlib import Path

from celery import Task
from sqlalchemy.orm import Session

from app.core.celery_app import celery_app
from app.core.config import settings
from app.db.session import SessionLocal
from app.models.document import Document, DocumentStatus, ProcessingStage, ExtractionMethod
from app.workers.chunking_tasks import chunk_document

# Import des processors hybrides
from app.rag.document_processor import DocumentProcessor, process_document_with_ocr
from app.rag.ocr_processor import MistralOCRClient

logger = logging.getLogger(__name__)


# =============================================================================
# UTILITAIRE : NETTOYAGE DU TEXTE
# =============================================================================

def sanitize_text_for_postgres(text: str) -> str:
    """
    Nettoyer le texte pour le rendre compatible avec PostgreSQL.
    
    PostgreSQL ne supporte pas les caractères NULL (\x00 ou \u0000) dans les 
    colonnes TEXT et JSONB. Cette fonction les supprime.
    
    Args:
        text: Texte à nettoyer
        
    Returns:
        Texte nettoyé sans caractères NULL
    """
    if not text:
        return text
    
    # Supprimer les caractères NULL (plusieurs méthodes pour être sûr)
    # Méthode 1: Remplacer \x00
    cleaned = text.replace('\x00', '')
    
    # Méthode 2: Supprimer via regex tous les caractères de contrôle problématiques
    # Garde les retours à la ligne (\n), tabulations (\t), et retours chariot (\r)
    cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', cleaned)
    
    return cleaned


def sanitize_dict_for_postgres(data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Nettoyer récursivement un dictionnaire pour PostgreSQL.
    
    Args:
        data: Dictionnaire à nettoyer
        
    Returns:
        Dictionnaire nettoyé
    """
    if not data:
        return data
    
    cleaned = {}
    for key, value in data.items():
        if isinstance(value, str):
            cleaned[key] = sanitize_text_for_postgres(value)
        elif isinstance(value, dict):
            cleaned[key] = sanitize_dict_for_postgres(value)
        elif isinstance(value, list):
            cleaned[key] = [
                sanitize_text_for_postgres(v) if isinstance(v, str)
                else sanitize_dict_for_postgres(v) if isinstance(v, dict)
                else v
                for v in value
            ]
        else:
            cleaned[key] = value
    
    return cleaned


# =============================================================================
# TASK BASE CLASS
# =============================================================================

class DocumentProcessingTask(Task):
    """Base task with database session management."""
    
    def __call__(self, *args, **kwargs):
        """Execute task with automatic session management."""
        return self.run(*args, **kwargs)
    
    def on_failure(self, exc, task_id, args, kwargs, einfo):
        """Handle task failure."""
        logger.error(f"Task {task_id} failed: {exc}")
        document_id = args[0] if args else kwargs.get('document_id')
        if document_id:
            db = SessionLocal()
            try:
                document = db.query(Document).filter(Document.id == document_id).first()
                if document:
                    document.status = DocumentStatus.FAILED
                    document.processing_stage = ProcessingStage.EXTRACTION
                    # Nettoyer le message d'erreur aussi
                    document.error_message = sanitize_text_for_postgres(str(exc))[:1000]
                    document.retry_count += 1
                    db.commit()
                    logger.info(f"Document {document_id} marked as FAILED")
            except Exception as e:
                logger.error(f"Failed to update document status: {e}")
                db.rollback()
            finally:
                db.close()


# =============================================================================
# PROCESSORS SINGLETON
# =============================================================================

_ocr_client = None
_document_processor = None


def get_ocr_client() -> Optional[MistralOCRClient]:
    """Get or create Mistral OCR client."""
    global _ocr_client
    if _ocr_client is None:
        if hasattr(settings, 'MISTRAL_API_KEY') and settings.MISTRAL_API_KEY:
            _ocr_client = MistralOCRClient(api_key=settings.MISTRAL_API_KEY)
            logger.info("Mistral OCR client initialized")
        else:
            logger.warning("MISTRAL_API_KEY not configured, OCR disabled")
    return _ocr_client


def get_document_processor() -> DocumentProcessor:
    """Get or create DocumentProcessor with OCR support."""
    global _document_processor
    if _document_processor is None:
        ocr_client = get_ocr_client()
        _document_processor = DocumentProcessor(ocr_client=ocr_client)
        logger.info("DocumentProcessor initialized with OCR support")
    return _document_processor


def _normalize_extraction_method(method: str) -> str:
    """Normaliser la méthode d'extraction en MAJUSCULES."""
    method_upper = method.upper()
    valid_methods = ['TEXT', 'OCR', 'HYBRID', 'FALLBACK']
    return method_upper if method_upper in valid_methods else 'TEXT'


# =============================================================================
# MAIN TASK
# =============================================================================

@celery_app.task(
    base=DocumentProcessingTask,
    bind=True,
    name="app.workers.processing_tasks.extract_document_text",
    max_retries=3,
    default_retry_delay=60,
    autoretry_for=(Exception,),
    retry_backoff=True,
    retry_backoff_max=600,
    retry_jitter=True
)
def extract_document_text(self, document_id: str) -> Dict[str, Any]:
    """
    Extract text from document using Hybrid DocumentProcessor.
    
    CORRECTIONS :
    - Nettoyage des caractères NULL avant stockage
    - Estimation correcte du nombre de pages
    - Enums en MAJUSCULES
    
    Args:
        document_id: UUID du document
        
    Returns:
        Dict avec toutes les métriques d'extraction
    """
    db = SessionLocal()
    start_time = datetime.utcnow()
    
    try:
        # 1. Charger le document
        document = db.query(Document).filter(Document.id == document_id).first()
        if not document:
            raise ValueError(f"Document {document_id} not found")
        
        logger.info(f"Starting extraction for document {document_id}: {document.original_filename}")
        
        # 2. Mettre à jour le statut
        document.status = DocumentStatus.PROCESSING
        document.processing_stage = ProcessingStage.EXTRACTION
        db.commit()
        
        # 3. Vérifier que le fichier existe
        file_path = Path(document.file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {document.file_path}")
        
        # 4. Utiliser DocumentProcessor HYBRIDE pour extraction
        processor = get_document_processor()
        
        try:
            # Extraction hybride
            result = processor.process_document(file_path)
            
            # CORRECTION : Nettoyer le texte extrait
            extracted_text = sanitize_text_for_postgres(result["text"])
            page_count = result.get("page_count", len(result.get("pages", [1])))
            total_images = result["total_images"]
            extraction_method = _normalize_extraction_method(result["extraction_method"])
            
            logger.info(
                f"DocumentProcessor extracted {len(extracted_text)} chars "
                f"from {page_count} pages, {total_images} images, "
                f"method={extraction_method} for {document.original_filename}"
            )
            
        except Exception as e:
            logger.error(f"DocumentProcessor failed, falling back to basic extraction: {e}")
            # Fallback vers extraction basique
            extracted_text, page_count = _fallback_extraction(file_path, document.file_extension)
            # CORRECTION : Nettoyer aussi le fallback
            extracted_text = sanitize_text_for_postgres(extracted_text)
            total_images = 0
            extraction_method = "FALLBACK"
        
        # 5. Calculer le temps d'extraction
        extraction_time = (datetime.utcnow() - start_time).total_seconds()
        
        # 6. MISE À JOUR DU DOCUMENT
        
        document.extracted_text_length = len(extracted_text)
        document.total_pages = page_count if page_count > 0 else 1
        document.extraction_time_seconds = extraction_time
        document.processing_stage = ProcessingStage.EXTRACTION
        
        # Colonnes OCR
        document.has_images = total_images > 0
        document.image_count = total_images
        document.extraction_method = extraction_method
        document.ocr_completed = extraction_method in ["OCR", "HYBRID"]
        
        # CORRECTION : Nettoyer les métadonnées JSONB
        metadata = {
            'extracted_text': extracted_text,
            'extraction_method': extraction_method,
            'total_images_ocr': total_images,
            'has_images': total_images > 0,
            'extraction_time_seconds': extraction_time,
            'pages_processed': page_count,
            'pages_estimated': document.file_extension.lower() in ['docx', 'doc', 'txt', 'md', 'rtf']
        }
        
        # Nettoyer tout le dictionnaire
        document.document_metadata = sanitize_dict_for_postgres(metadata)
        
        db.commit()
        
        logger.info(
            f"Extraction completed for {document_id}: "
            f"{len(extracted_text)} chars, {page_count} pages, "
            f"{total_images} images OCR, method={extraction_method}, "
            f"{extraction_time:.2f}s"
        )
        
        # 7. Envoyer vers la queue chunking
        chunk_document.apply_async(args=[str(document_id)], queue='chunking')
        
        return {
            "document_id": str(document_id),
            "extracted_text_length": len(extracted_text),
            "page_count": page_count,
            "pages_estimated": document.file_extension.lower() in ['docx', 'doc', 'txt', 'md', 'rtf'],
            "extraction_time_seconds": extraction_time,
            "extraction_method": extraction_method,
            "total_images_ocr": total_images,
            "has_images": total_images > 0,
            "ocr_completed": extraction_method in ["OCR", "HYBRID"]
        }
        
    except Exception as e:
        logger.error(f"Extraction failed for document {document_id}: {e}")
        db.rollback()
        
        # Mettre à jour l'erreur
        document = db.query(Document).filter(Document.id == document_id).first()
        if document:
            # CORRECTION : Nettoyer le message d'erreur
            document.error_message = sanitize_text_for_postgres(str(e))[:1000]
            document.retry_count += 1
            db.commit()
        
        raise
        
    finally:
        db.close()


# =============================================================================
# FALLBACK EXTRACTION
# =============================================================================

CHARS_PER_PAGE_ESTIMATE = 2500


def _estimate_page_count(text_length: int) -> int:
    """Estimer le nombre de pages à partir de la longueur du texte."""
    import math
    pages = math.ceil(text_length / CHARS_PER_PAGE_ESTIMATE)
    return max(1, pages)


def _fallback_extraction(file_path: Path, file_extension: str) -> tuple[str, int]:
    """Fallback extraction basique si DocumentProcessor échoue."""
    file_extension = file_extension.lower().lstrip('.')
    
    if file_extension == 'pdf':
        return _extract_from_pdf(file_path)
    elif file_extension in ['docx', 'doc']:
        text = _extract_from_docx(file_path)
        return text, _estimate_page_count(len(text))
    elif file_extension in ['xlsx', 'xls']:
        text, sheet_count = _extract_from_xlsx(file_path)
        return text, sheet_count
    elif file_extension in ['pptx', 'ppt']:
        return _extract_from_pptx(file_path)
    elif file_extension == 'txt':
        text = _extract_from_txt(file_path)
        return text, _estimate_page_count(len(text))
    elif file_extension == 'md':
        text = _extract_from_markdown(file_path)
        return text, _estimate_page_count(len(text))
    elif file_extension == 'rtf':
        text = _extract_from_rtf(file_path)
        return text, _estimate_page_count(len(text))
    elif file_extension in ['png', 'jpg', 'jpeg', 'webp', 'gif']:
        return "[Image - OCR required but unavailable]", 1
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")


def _extract_from_pdf(file_path: Path) -> tuple[str, int]:
    """Extraction PDF basique."""
    try:
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader
        
        extracted_text = ""
        reader = PdfReader(str(file_path))
        page_count = len(reader.pages)
        
        for page_num, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                if text:
                    # CORRECTION : Nettoyer chaque page
                    text = sanitize_text_for_postgres(text)
                    extracted_text += f"\n\n--- Page {page_num} ---\n\n{text}"
            except Exception as e:
                logger.warning(f"Failed to extract text from page {page_num}: {e}")
        
        return extracted_text, page_count
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        raise ValueError(f"Failed to extract PDF: {e}")


def _extract_from_docx(file_path: Path) -> str:
    """Extraction DOCX basique."""
    try:
        from docx import Document as DocxDocument
        
        doc = DocxDocument(str(file_path))
        extracted_text = ""
        
        for para in doc.paragraphs:
            if para.text.strip():
                extracted_text += para.text + "\n\n"
        
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                extracted_text += row_text + "\n"
        
        return extracted_text
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        raise ValueError(f"Failed to extract DOCX: {e}")


def _extract_from_xlsx(file_path: Path) -> tuple[str, int]:
    """Extraction XLSX basique."""
    try:
        from openpyxl import load_workbook
        
        workbook = load_workbook(str(file_path), data_only=True)
        extracted_text = ""
        sheet_count = len(workbook.sheetnames)
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            extracted_text += f"\n\n[Feuille: {sheet_name}]\n\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    extracted_text += row_text + "\n"
        
        return extracted_text, sheet_count
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        raise ValueError(f"Failed to extract XLSX: {e}")


def _extract_from_pptx(file_path: Path) -> tuple[str, int]:
    """Extraction PPTX basique."""
    try:
        from pptx import Presentation
        
        prs = Presentation(str(file_path))
        extracted_text = ""
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            extracted_text += f"\n\n--- Slide {slide_num} ---\n\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    extracted_text += shape.text + "\n"
                
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        extracted_text += row_text + "\n"
        
        return extracted_text, slide_count
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        raise ValueError(f"Failed to extract PPTX: {e}")


def _extract_from_txt(file_path: Path) -> str:
    """Extraction TXT basique."""
    encodings = ['utf-8', 'latin-1', 'cp1252']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            continue
    raise ValueError("Failed to decode text file")


def _extract_from_markdown(file_path: Path) -> str:
    """Extraction Markdown basique."""
    return _extract_from_txt(file_path)


def _extract_from_rtf(file_path: Path) -> str:
    """Extraction RTF basique."""
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return rtf_to_text(f.read())
    except ImportError:
        import re
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        content = re.sub(r'\\[a-z]+\d*\s?', '', content)
        return re.sub(r'[{}]', '', content)